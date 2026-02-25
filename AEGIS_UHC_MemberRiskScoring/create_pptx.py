from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = r"D:\AI_Products_Dev\Products\AgenticAIClass_Docs\AEGIS_MemberRiskScoring_VP_CFO_Deck_v3.pptx"

NAVY     = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x44, 0x72, 0xC4)
STEEL    = RGBColor(0x2E, 0x75, 0xB6)
DARK_SL  = RGBColor(0x2F, 0x54, 0x96)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x33, 0x33, 0x33)
ACCENT   = RGBColor(0xD2, 0x60, 0x12)
LT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
LIGHT_BG = RGBColor(0xEB, 0xF1, 0xF8)
GRAY     = RGBColor(0x88, 0x88, 0x88)
GREEN    = RGBColor(0x00, 0x80, 0x00)
RED      = RGBColor(0xCC, 0x00, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_rect(slide, x, y, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, x, y, w, h, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, x, y, w, h, items, size=13, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_header_bar(slide):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), NAVY)

def add_footer(slide, text="CONFIDENTIAL  |  UHC AEGIS Platform  |  February 2026"):
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
                 text, size=8, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(sl, Inches(0), Inches(5.8), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "AEGIS", size=54, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Advanced Early Warning Signal Intelligent System", size=26, color=STEEL, bold=False, align=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Member Risk Scoring  &  Provider Risk Scoring", size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
             "Executive Briefing  |  VP & CFO Review  |  February 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "UnitedHealth Group  |  Data Science & Analytics", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — WHY NOW: THE MLR CRISIS (Data-Driven Story)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
             "Why Now?  The Medical Cost Crisis Is Accelerating", size=26, color=NAVY, bold=True)

add_text_box(sl, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
             '"The full year 2025 adjusted medical care ratio was 88.9% compared to 85.5% in 2024 ... CMS Medicare funding reductions\nand changes from the Inflation Reduction Act combined with accelerating medical cost trends."',
             size=11, color=GRAY, bold=False)

SRC_CLR = RGBColor(0x55, 0x77, 0x99)
add_text_box(sl, Inches(0.6), Inches(1.45), Inches(12), Inches(0.3),
             "Source: UNH Q4 2025 Earnings Release, Jan 27, 2026  |  unitedhealthgroup.com/investors/financial-reports",
             size=8, color=SRC_CLR, bold=False)

# VISUAL: MCR Bar Chart (simulated with rectangles)
add_text_box(sl, Inches(0.6), Inches(1.85), Inches(5.5), Inches(0.4),
             "Medical Care Ratio  (Annual)", size=15, color=NAVY, bold=True)

BAR_RED   = RGBColor(0xC0, 0x39, 0x2B)
BAR_AMBER = RGBColor(0xD4, 0x8B, 0x2D)
BAR_GREEN = RGBColor(0x27, 0xAE, 0x60)

mcr_data = [
    ("2023", 83.2, BAR_GREEN, "$241.9B", "$290.8B"),
    ("2024", 85.5, BAR_AMBER, "$264.2B", "$308.8B"),
    ("2025", 89.1, BAR_RED,   "$314.0B", "$352.2B"),
    ("2026E", 88.8, MED_BLUE, "Guidance", ">$439B rev"),
]
bar_base_y = Inches(5.7)
bar_max_h = Inches(3.2)
mcr_min = 78.0
mcr_range = 14.0

for i, (year, mcr, clr, med_cost, prem) in enumerate(mcr_data):
    bx = Inches(1.0 + i * 1.35)
    bar_h_ratio = (mcr - mcr_min) / mcr_range
    bar_h = bar_max_h * bar_h_ratio
    bar_y = bar_base_y - bar_h

    is_2026 = year == "2026E"
    if is_2026:
        r = add_rect(sl, bx, bar_y, Inches(1.1), bar_h, clr, clr)
        r.fill.background()
        r.line.color.rgb = clr
        r.line.width = Pt(2)
    else:
        add_rect(sl, bx, bar_y, Inches(1.1), bar_h, clr)

    add_text_box(sl, bx, bar_y - Inches(0.35), Inches(1.1), Inches(0.35),
                 f"{mcr}%", size=16, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, bx, Inches(5.75), Inches(1.1), Inches(0.3),
                 year, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# Delta annotations
add_text_box(sl, Inches(2.1), Inches(2.3), Inches(1.2), Inches(0.3),
             "+2.3 pts", size=10, color=BAR_AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, Inches(3.5), Inches(2.15), Inches(1.2), Inches(0.3),
             "+3.6 pts", size=10, color=BAR_RED, bold=True, align=PP_ALIGN.CENTER)

# RIGHT SIDE: The Financial Impact
add_text_box(sl, Inches(6.8), Inches(1.85), Inches(6), Inches(0.4),
             "What The Earnings Reports Reveal", size=15, color=NAVY, bold=True)

impact_cards = [
    ("360 bps", "MCR Deterioration\nin Just 2 Years",
     '"The full year 2025 adjusted medical care ratio was 88.9% compared to 85.5% in 2024." Medical costs grew faster than premium revenue -- the pricing gap widened every quarter.',
     BAR_RED),
    ("2.7%", "Operating Margin\nDown from 5.2%",
     '"Operating margin of 2.7% compared to 5.2% in 2024, primarily due to Biden-era Medicare funding reductions and Inflation Reduction Act impacts combined with elevated medical cost trends."',
     BAR_AMBER),
    ("92.4%", "Q4 2025 MCR\nHighest in a Decade",
     'UNH took a restructuring charge in Q4 and restructured leadership. "The company laid the foundation for more disciplined and transparent operations, stronger performance and sustained growth."',
     NAVY),
]

for i, (big_num, subtitle, detail, clr) in enumerate(impact_cards):
    cy = Inches(2.4 + i * 1.55)
    add_rect(sl, Inches(6.8), cy, Inches(1.6), Inches(1.35), clr)
    add_text_box(sl, Inches(6.85), cy + Inches(0.1), Inches(1.5), Inches(0.55),
                 big_num, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(6.85), cy + Inches(0.65), Inches(1.5), Inches(0.6),
                 subtitle, size=8, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(8.55), cy + Inches(0.05), Inches(4.2), Inches(1.25),
                 detail, size=10, color=BLACK)

# Source footer
add_text_box(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
             "Sources:  [1] UNH 2025 Annual Earnings Release (Jan 27, 2026) - businesswire.com/news/home/20260126830491\n"
             "[2] UNH Q4 2024 Earnings Release (Jan 16, 2025) - unitedhealthgroup.com/investors  |  [3] Becker's Payer Issues, Jan 2026",
             size=7, color=SRC_CLR)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — QUARTERLY MLR TREND + ACTUARIAL PRICING GAP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "The Pricing Gap  —  Why Actuaries Need Early Signals", size=26, color=NAVY, bold=True)

add_text_box(sl, Inches(0.6), Inches(0.9), Inches(12), Inches(0.5),
             "Premiums are priced using backward-looking assumptions. When medical cost trends accelerate faster than pricing can adjust, the gap widens every quarter.",
             size=13, color=BLACK)

# Quarterly MCR table
add_text_box(sl, Inches(0.6), Inches(1.55), Inches(5), Inches(0.4),
             "Quarterly MCR Trajectory", size=15, color=NAVY, bold=True)

q_headers = ["Quarter", "2024 MCR", "2025 MCR", "YoY Gap"]
q_data = [
    ("Q1", "84.3%", "84.8%", "+0.5 pts"),
    ("Q2", "85.1%", "89.4%", "+4.3 pts"),
    ("Q3", "85.2%", "89.9%", "+4.7 pts"),
    ("Q4", "~87.6%", "92.4%", "+4.8 pts"),
    ("Full Year", "85.5%", "89.1%", "+3.6 pts"),
]

ty = Inches(2.05)
col_xs = [Inches(0.6), Inches(2.3), Inches(4.0), Inches(5.7)]
col_ws = [Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.8)]

for j, hdr in enumerate(q_headers):
    add_rect(sl, col_xs[j], ty, col_ws[j], Inches(0.4), NAVY)
    add_text_box(sl, col_xs[j] + Inches(0.05), ty + Inches(0.02), col_ws[j] - Inches(0.1), Inches(0.36),
                 hdr, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for ri, row in enumerate(q_data):
    ry = ty + Inches(0.42) + Inches(ri * 0.38)
    is_total = ri == len(q_data) - 1
    row_bg = LIGHT_BG if ri % 2 == 0 else WHITE
    if is_total:
        row_bg = RGBColor(0xE8, 0xD5, 0xC4)

    for j, val in enumerate(row):
        add_rect(sl, col_xs[j], ry, col_ws[j], Inches(0.36), row_bg, STEEL if is_total else None)
        font_clr = NAVY if is_total else (BAR_RED if j == 3 and not is_total and ri > 0 else BLACK)
        add_text_box(sl, col_xs[j] + Inches(0.05), ry + Inches(0.01), col_ws[j] - Inches(0.1), Inches(0.34),
                     val, size=10 if not is_total else 11,
                     color=font_clr, bold=is_total, align=PP_ALIGN.CENTER)

# RIGHT: The Story - "What went wrong"
add_text_box(sl, Inches(9.2), Inches(1.55), Inches(3.8), Inches(0.4),
             "The Accelerating Gap", size=15, color=NAVY, bold=True)

gap_story = [
    ("Q1 2025: Manageable", "+0.5 pts", "Still within actuarial tolerance. Pricing adjustments proceeding normally.",
     BAR_GREEN),
    ("Q2 2025: Warning Signs", "+4.3 pts", "Specialty Rx costs spiked. Hospital coding intensity increased. IRA Part D changes hit.",
     BAR_AMBER),
    ("Q3 2025: Crisis Level", "+4.7 pts", "Medical trend outpacing pricing. UNH raised earnings outlook but MCR kept climbing quarter after quarter.",
     BAR_RED),
    ("Q4 2025: Emergency", "+4.8 pts", "MCR hit 92.4%. Restructuring charge taken. Company restructured leadership and operations.",
     RGBColor(0x8B, 0x00, 0x00)),
]

for i, (phase, delta, desc, clr) in enumerate(gap_story):
    gy = Inches(2.05 + i * 1.15)
    add_rect(sl, Inches(9.2), gy, Inches(0.15), Inches(0.95), clr)
    add_text_box(sl, Inches(9.5), gy, Inches(3.3), Inches(0.3),
                 phase, size=11, color=clr, bold=True)
    add_text_box(sl, Inches(9.5), gy + Inches(0.3), Inches(3.3), Inches(0.6),
                 desc, size=9, color=BLACK)

# Bottom: THE CORE INSIGHT
add_rect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), NAVY)
add_text_box(sl, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.5),
             "THE CORE OBJECTIVE:  Predict Member-Level Medical Spend Before Actuaries Set Prices",
             size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

obj_bullets = [
    "If we can predict which members will become high-cost NEXT QUARTER, actuaries can price premiums with 90-day forward visibility instead of 12-month backward assumptions",
    "Every 100 basis points of MCR improvement directly recovers operating margin -- the difference between 2.7% and 5.2% margin that UHC lost in one year",
    "Early signals feed three critical functions:  (1) Care Management for intervention  (2) Actuarial for pricing accuracy  (3) Stop-Loss for reserve adequacy",
]
for i, b in enumerate(obj_bullets):
    add_text_box(sl, Inches(1.0), Inches(5.85 + i * 0.35), Inches(11.5), Inches(0.35),
                 f"\u25b8  {b}", size=10, color=LT_BLUE)

add_text_box(sl, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             "Sources:  UNH Q1-Q4 2025 Earnings Releases - unitedhealthgroup.com/investors  |  FierceHealthcare: \"UHC expects to lose 1.3M MA members\" Jan 2026\n"
             "Norwood: \"UHC's MLR an eye-opening reminder of V28 impact\" - norwood.com  |  Becker's Payer Issues  |  SEC filings: 10-K, 10-Q (2024-2025)",
             size=7, color=SRC_CLR)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Three Compounding Problems", size=28, color=NAVY, bold=True)

add_text_box(sl, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
             "The MLR crisis is a symptom. These are the root causes we can address with data science.",
             size=14, color=BLACK)

boxes = [
    ("High-Cost Claimant\nConcentration",
     "2.2% of members generate 23% of total medical spend. These members ARE identifiable 90 days before cost spikes, but we have no systematic prediction engine.\n\nWith 89.1% MCR, every avoidable high-cost event directly erodes operating margin.",
     "2.2% of members = 23% of spend\nPredictable 90 days early"),
    ("Actuarial Pricing\nBlind Spot",
     "Premiums are set using population-level historical trends. But medical costs accelerated faster than pricing could adjust -- MCR jumped 360 bps in 2 years.\n\nMember-level spend prediction gives Actuarial a 90-day forward view to price with, not a 12-month backward mirror.",
     "MCR went from 83.2% to 89.1%\nPricing could not keep up"),
    ("No Unified\nEarly Warning System",
     "Care Management, Actuarial, Stop-Loss, and Network teams all react AFTER costs are incurred. No shared predictive signal connects them.\n\nQ4 2025 MCR hit 92.4% -- proving that by the time lagging indicators catch up, the damage is done.",
     "Q4 2025 MCR: 92.4%\nReactive = too late"),
]

for i, (title, desc, stat) in enumerate(boxes):
    bx = Inches(0.6 + i * 4.1)
    r = add_rect(sl, bx, Inches(1.7), Inches(3.8), Inches(4.8), LIGHT_BG, STEEL)
    add_text_box(sl, bx + Inches(0.2), Inches(1.85), Inches(3.4), Inches(0.6),
                 title, size=15, color=NAVY, bold=True)
    add_text_box(sl, bx + Inches(0.2), Inches(2.55), Inches(3.4), Inches(2.8),
                 desc, size=11, color=BLACK)
    r2 = add_rect(sl, bx + Inches(0.15), Inches(5.4), Inches(3.5), Inches(0.9), NAVY)
    add_text_box(sl, bx + Inches(0.25), Inches(5.45), Inches(3.3), Inches(0.8),
                 stat, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — OVERALL OBJECTIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Overall Objective", size=28, color=NAVY, bold=True)

add_text_box(sl, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             'Build "AEGIS" — a unified early warning platform that predicts member cost risk AND provider performance risk, presented through a single Executive Dashboard for proactive decision-making.',
             size=15, color=BLACK, bold=False)

# Two columns
col_data = [
    ("Member Risk Scoring", NAVY, [
        "Predict which members will become high-cost claimants next quarter",
        "Estimate expected spend per member with confidence intervals",
        "Assign risk tiers: CRITICAL | HIGH | ELEVATED | WATCH",
        "Provide SHAP-based explainability (top cost drivers per member)",
        "Feed Care Management, Actuarial, and Stop-Loss teams weekly",
    ]),
    ("Provider Risk Scoring", DARK_SL, [
        "Cluster providers by geo, specialty, and practice embeddings",
        "Compute Economic Performance Index (denial ratio, 1st-pass rate, cost efficiency)",
        "Compute Administrative / Coding Maturity Index",
        "Link member risk scores to provider panels (cross-input)",
        "Feed Network Management, Contracting, and Value-Based teams",
    ]),
]

for i, (title, bg, items) in enumerate(col_data):
    cx = Inches(0.6 + i * 6.2)
    add_rect(sl, cx, Inches(2.2), Inches(5.8), Inches(0.5), bg)
    add_text_box(sl, cx + Inches(0.15), Inches(2.22), Inches(5.5), Inches(0.45),
                 title, size=16, color=WHITE, bold=True)
    add_bullet_list(sl, cx + Inches(0.2), Inches(2.9), Inches(5.5), Inches(3.5),
                    [f"\u2022  {it}" for it in items], size=12, color=BLACK)

# Bottom bar
add_rect(sl, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5), ACCENT)
add_text_box(sl, Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.45),
             "Both streams converge into a single Early Warning Intelligence Signal Dashboard",
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — MEMBER RISK SCORING DEEP DIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Member Risk Scoring  —  Approach & Model", size=28, color=NAVY, bold=True)

# Left column
add_text_box(sl, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
             "Two-Stage ML Pipeline", size=16, color=NAVY, bold=True)
add_bullet_list(sl, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5), [
    "\u2022  Stage 1: XGBoost Binary Classifier",
    "     \"Will this member become high-cost next quarter?\"",
    "     85 features across 6 categories",
    "     Target: AUC-ROC > 0.85",
    "",
    "\u2022  Stage 2: LightGBM Cost Regression",
    "     \"HOW MUCH will they cost?\"",
    "     Log-transformed target for skewed spend",
    "     Confidence intervals per member",
], size=12, color=BLACK)

# Right column - Risk Tiers
add_text_box(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(0.4),
             "Risk Tier Output", size=16, color=NAVY, bold=True)

tiers = [
    ("CRITICAL", "> 0.70", "Top 0.5%", "Immediate CM outreach", ACCENT),
    ("HIGH", "0.50-0.70", "Top 1.5%", "Proactive care plan", RGBColor(0xD4, 0x8B, 0x2D)),
    ("ELEVATED", "0.30-0.50", "Top 4%", "Monitor + targeted programs", MED_BLUE),
    ("WATCH", "0.15-0.30", "Top 9%", "Quarterly reassessment", STEEL),
    ("STANDARD", "< 0.15", "~91%", "Standard population health", GRAY),
]

ty = Inches(1.75)
for name, prob, pctl, action, clr in tiers:
    r = add_rect(sl, Inches(7), ty, Inches(1.5), Inches(0.42), clr)
    add_text_box(sl, Inches(7.1), ty + Inches(0.02), Inches(1.3), Inches(0.38),
                 name, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(8.6), ty + Inches(0.02), Inches(1.3), Inches(0.38),
                 f"P: {prob}", size=10, color=BLACK)
    add_text_box(sl, Inches(9.9), ty + Inches(0.02), Inches(1.0), Inches(0.38),
                 pctl, size=10, color=BLACK)
    add_text_box(sl, Inches(10.9), ty + Inches(0.02), Inches(1.8), Inches(0.38),
                 action, size=9, color=GRAY)
    ty += Inches(0.48)

# Bottom - Key numbers
add_rect(sl, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.06), NAVY)
add_text_box(sl, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
             "Scale & Impact", size=16, color=NAVY, bold=True)

metrics = [
    ("50M+", "Total Members"),
    ("2.2%", "High-Cost Concentration"),
    ("23%", "of Spend from Top 2.2%"),
    ("85", "ML Features"),
    ("90 days", "Prediction Horizon"),
    ("4 Weeks", "MVP Timeline"),
]
for i, (val, label) in enumerate(metrics):
    mx = Inches(0.6 + i * 2.05)
    add_rect(sl, mx, Inches(5.5), Inches(1.85), Inches(1.2), LIGHT_BG, STEEL)
    add_text_box(sl, mx, Inches(5.55), Inches(1.85), Inches(0.6),
                 val, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, mx, Inches(6.15), Inches(1.85), Inches(0.4),
                 label, size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — PROVIDER RISK SCORING DEEP DIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Provider Risk Scoring  —  Approach & Indices", size=28, color=NAVY, bold=True)

# Three columns for the three components
components = [
    ("Provider Clustering", NAVY, [
        "Group 1.2M+ providers into peer cohorts",
        "Features: geo (lat/lon), specialty, payer mix, case mix, volume",
        "Method: UMAP + HDBSCAN embeddings",
        "Output: Peer groups for fair comparison",
        "Example: Urban Cardiologists, Rural PCPs, Academic Oncologists",
    ]),
    ("Admin / Coding\nMaturity Index", DARK_SL, [
        "First-Pass Claim Rate (30% weight)",
        "Denial Rate inverse (25% weight)",
        "Coding Specificity Score (20% weight)",
        "Avg Days to Submit inverse (15% weight)",
        "Rework Rate inverse (10% weight)",
        "Output: 0-100 composite score",
    ]),
    ("Economic\nPerformance Index", ACCENT, [
        "Risk-Adjusted Cost Index / O:E (35%)",
        "First-Pass Rate (15%)",
        "Denial Rate inverse (15%)",
        "Network Leakage inverse (10%)",
        "Readmission Rate inverse (10%)",
        "Cost Trend inverse (10%)",
        "Total Cost of Care PMPY (5%)",
    ]),
]

for i, (title, bg, items) in enumerate(components):
    cx = Inches(0.5 + i * 4.2)
    add_rect(sl, cx, Inches(1.2), Inches(3.9), Inches(0.65), bg)
    add_text_box(sl, cx + Inches(0.1), Inches(1.22), Inches(3.7), Inches(0.6),
                 title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(sl, cx + Inches(0.15), Inches(2.0), Inches(3.6), Inches(3.5),
                    [f"\u2022  {it}" for it in items], size=11, color=BLACK, spacing=Pt(4))

# Cross-input callout
add_rect(sl, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.55), LIGHT_BG, STEEL)
add_text_box(sl, Inches(0.7), Inches(5.62), Inches(11.9), Inches(0.5),
             "Cross-Input:  Member Risk Scores flow INTO Provider Scoring  \u2192  % of panel flagged CRITICAL/HIGH + avg member risk score per provider",
             size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# Provider tiers
add_text_box(sl, Inches(0.6), Inches(6.3), Inches(3), Inches(0.4),
             "Provider Risk Tiers:", size=14, color=NAVY, bold=True)
ptiers = [
    ("TOP PERFORMER", "> 80", GREEN),
    ("STANDARD", "60-80", MED_BLUE),
    ("WATCH", "40-60", RGBColor(0xD4, 0x8B, 0x2D)),
    ("CONCERN", "< 40", RED),
]
for i, (name, score, clr) in enumerate(ptiers):
    px = Inches(3.8 + i * 2.4)
    add_rect(sl, px, Inches(6.3), Inches(2.1), Inches(0.42), clr)
    add_text_box(sl, px + Inches(0.05), Inches(6.32), Inches(2.0), Inches(0.38),
                 f"{name}  ({score})", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "End-to-End Architecture", size=28, color=NAVY, bold=True)

# Member box
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6.0), Inches(4.0), LIGHT_BG, STEEL)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.45), NAVY)
add_text_box(sl, Inches(0.5), Inches(1.22), Inches(5.8), Inches(0.42),
             "MEMBER RISK SCORING", size=14, color=WHITE, bold=True)

m_sources = ["Medical Claims\n837P/837I", "Call Center\nIVR/CRM", "Remittance\n835/EOB", "Enrollment\n834"]
for i, src in enumerate(m_sources):
    add_rect(sl, Inches(0.6 + i * 1.4), Inches(1.85), Inches(1.25), Inches(0.7), NAVY)
    add_text_box(sl, Inches(0.65 + i * 1.4), Inches(1.87), Inches(1.15), Inches(0.65),
                 src, size=8, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.7), Inches(2.8), Inches(2.3), Inches(0.45), MED_BLUE)
add_text_box(sl, Inches(0.75), Inches(2.82), Inches(2.2), Inches(0.4),
             "ETL \u2192 Feature Store (85)", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(3.3), Inches(2.8), Inches(2.8), Inches(0.45), ACCENT)
add_text_box(sl, Inches(3.35), Inches(2.82), Inches(2.7), Inches(0.4),
             "XGBoost \u2192 LightGBM", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(1.5), Inches(3.6), Inches(3.8), Inches(0.65), NAVY)
add_text_box(sl, Inches(1.55), Inches(3.62), Inches(3.7), Inches(0.6),
             "Member Risk Score Table", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Provider box
add_rect(sl, Inches(6.9), Inches(1.2), Inches(6.0), Inches(4.0), RGBColor(0xF5, 0xF5, 0xF0), STEEL)
add_rect(sl, Inches(6.9), Inches(1.2), Inches(6.0), Inches(0.45), DARK_SL)
add_text_box(sl, Inches(7.0), Inches(1.22), Inches(5.8), Inches(0.42),
             "PROVIDER RISK SCORING", size=14, color=WHITE, bold=True)

p_sources = ["Provider NPI\nNPPES", "Claims per\nProvider", "New Drug\nFormulary", "External\nGeo/Market"]
for i, src in enumerate(p_sources):
    add_rect(sl, Inches(7.1 + i * 1.4), Inches(1.85), Inches(1.25), Inches(0.7), NAVY)
    add_text_box(sl, Inches(7.15 + i * 1.4), Inches(1.87), Inches(1.15), Inches(0.65),
                 src, size=8, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(7.2), Inches(2.8), Inches(2.3), Inches(0.45), MED_BLUE)
add_text_box(sl, Inches(7.25), Inches(2.82), Inches(2.2), Inches(0.4),
             "Clustering + Indices", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(9.8), Inches(2.8), Inches(2.8), Inches(0.45), ACCENT)
add_text_box(sl, Inches(9.85), Inches(2.82), Inches(2.7), Inches(0.4),
             "Econ Perf + Admin Idx", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(8.0), Inches(3.6), Inches(3.8), Inches(0.65), DARK_SL)
add_text_box(sl, Inches(8.05), Inches(3.62), Inches(3.7), Inches(0.6),
             "Provider Risk Score Table", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Cross-input arrow label
add_text_box(sl, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.3),
             "+ Member Scores \u2192", size=9, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Dashboard
add_rect(sl, Inches(2.5), Inches(5.6), Inches(8.3), Inches(0.65), NAVY)
add_text_box(sl, Inches(2.55), Inches(5.62), Inches(8.2), Inches(0.6),
             "Early Warning Intelligence Signal Dashboard", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# End users
users = ["Care Mgmt", "Actuarial", "Stop-Loss", "Provider\nNetwork", "Executive", "Underwriting", "Contracting"]
for i, u in enumerate(users):
    add_rect(sl, Inches(0.5 + i * 1.8), Inches(6.6), Inches(1.6), Inches(0.55), STEEL)
    add_text_box(sl, Inches(0.55 + i * 1.8), Inches(6.62), Inches(1.5), Inches(0.5),
                 u, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — TIMELINE & ROADMAP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Timeline & Roadmap", size=28, color=NAVY, bold=True)

# Timeline bar
add_rect(sl, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.06), NAVY)

phases = [
    ("MVP 1\nMar 1, 2026", "Member Scoring\n10% of Members\n(~5M members)", ACCENT,
     "4 weeks\nFeb 3 - Mar 1",
     ["\u2022 Extract 24mo claims + enrollment",
      "\u2022 Build 85-feature matrix",
      "\u2022 Train XGBoost + LightGBM",
      "\u2022 SHAP explainability",
      "\u2022 Risk tier assignment",
      "\u2022 Dashboard v1 (member view)"]),
    ("MVP 2\nApr 1, 2026", "Member Scoring\n100% Scale-Up\n(50M members)", NAVY,
     "4 weeks\nMar 2 - Apr 1",
     ["\u2022 Scale pipeline to full population",
      "\u2022 Add SDOH features (5 free sources)",
      "\u2022 Fairness audit + calibration",
      "\u2022 Weekly CM feed automation",
      "\u2022 Actuarial / stop-loss integration",
      "\u2022 Production monitoring (PSI)"]),
    ("MVP 3\nMay 1, 2026", "Provider Scoring\nAll 1.2M Providers", DARK_SL,
     "4 weeks\nApr 2 - May 1",
     ["\u2022 Provider clustering (embeddings)",
      "\u2022 Admin/Coding Maturity Index",
      "\u2022 Economic Performance Index",
      "\u2022 Cross-link member scores",
      "\u2022 Provider risk tiers",
      "\u2022 Dashboard v2 (member + provider)"]),
    ("Full Platform\nJun 1, 2026", "Combined Dashboard\n+ VBC Integration", STEEL,
     "4 weeks\nMay 2 - Jun 1",
     ["\u2022 VBC contract integration",
      "\u2022 Provider outreach automation",
      "\u2022 Real-time scoring API",
      "\u2022 User training + rollout",
      "\u2022 Phase 2 roadmap kickoff",
      "\u2022 A/B testing framework"]),
]

for i, (title, scope, color, duration, tasks) in enumerate(phases):
    px = Inches(0.5 + i * 3.15)
    # Milestone marker
    add_rect(sl, px + Inches(1.2), Inches(1.15), Inches(0.2), Inches(0.35), color)
    add_text_box(sl, px, Inches(1.55), Inches(2.8), Inches(0.9),
                 title, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_rect(sl, px + Inches(0.1), Inches(2.5), Inches(2.8), Inches(0.85), color)
    add_text_box(sl, px + Inches(0.15), Inches(2.52), Inches(2.7), Inches(0.8),
                 scope, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(sl, px + Inches(0.1), Inches(3.45), Inches(2.8), Inches(0.3),
                 duration, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    add_bullet_list(sl, px + Inches(0.1), Inches(3.8), Inches(2.8), Inches(3.0),
                    tasks, size=10, color=BLACK, spacing=Pt(3))

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — EXPECTED IMPACT (Operational, No $ Assumptions)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Expected Operational Impact", size=28, color=NAVY, bold=True)

add_text_box(sl, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
             "AEGIS directly addresses the three drivers cited in UHC earnings releases: elevated medical cost trends, hospital coding intensity, and actuarial pricing gaps.",
             size=14, color=BLACK)

impact_cols = [
    ("MCR Improvement\nLevers", NAVY, [
        "Identify high-cost members 90 days before cost spikes",
        "Enable proactive Care Management intervention before claims incur",
        "Shift actuarial pricing from 12-month backward to 90-day forward view",
        "Improve IBNR reserve accuracy with member-level predictions",
        "Reduce stop-loss mispricing via better tail-risk visibility",
    ]),
    ("Operational\nEfficiency Gains", DARK_SL, [
        "Automate weekly risk-scored member lists to Care Management",
        "Replace manual provider performance reports with real-time indices",
        "Unify siloed signals (claims, Rx, SDOH) into one predictive score",
        "Enable data-driven provider contracting and VBC negotiations",
        "Reduce time-to-insight from months to days",
    ]),
    ("Strategic\nAdvantage", ACCENT, [
        "First-mover predictive platform across 50M+ members",
        "Provider risk scoring creates leverage for network negotiations",
        "SHAP explainability builds regulator and clinician trust",
        "Scalable framework: add new use cases (readmission, Rx adherence)",
        "Positions UHC to meet 2026 MCR guidance of 88.8%",
    ]),
]

for i, (title, bg, items) in enumerate(impact_cols):
    cx = Inches(0.5 + i * 4.2)
    add_rect(sl, cx, Inches(1.7), Inches(3.9), Inches(0.65), bg)
    add_text_box(sl, cx + Inches(0.1), Inches(1.72), Inches(3.7), Inches(0.6),
                 title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(sl, cx + Inches(0.15), Inches(2.5), Inches(3.6), Inches(3.5),
                    [f"\u2022  {it}" for it in items], size=11, color=BLACK, spacing=Pt(5))

# Bottom - The key metric
add_rect(sl, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), NAVY)
add_text_box(sl, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.35),
             "The Benchmark:  Every 100 basis points of MCR improvement directly restores operating margin",
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(sl, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.35),
             "UHC went from 5.2% to 2.7% operating margin in one year.  Predictive signals are the fastest path to reverse that trend.",
             size=12, color=LT_BLUE, align=PP_ALIGN.CENTER)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — ASK / NEXT STEPS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(sl)
add_text_box(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Ask & Next Steps", size=28, color=NAVY, bold=True)

asks = [
    ("1", "Approve MVP 1 Sprint Start", "Feb 3, 2026",
     "Green-light the 4-week sprint to build Member Risk Scoring on 10% of the member population. Deliver scored risk list to Care Management by March 1."),
    ("2", "Data Access Authorization", "Immediate",
     "Grant analytics team access to 24 months of claims (837P/I), enrollment (834), and remittance (835) data in the analytics environment."),
    ("3", "Stakeholder Alignment", "Week 1",
     "Schedule 30-min sessions with Care Management, Actuarial, and Network leads to align on risk tier thresholds and dashboard requirements."),
    ("4", "Fund Full Platform Build-Out", "By Apr 2026",
     "Based on MVP 1 results, approve Year 1 budget for full AEGIS platform: member + provider scoring, dashboard, production pipeline, and team allocation."),
]

for i, (num, title, when, desc) in enumerate(asks):
    ay = Inches(1.2 + i * 1.45)
    add_rect(sl, Inches(0.6), ay, Inches(0.6), Inches(1.2), NAVY)
    add_text_box(sl, Inches(0.65), ay + Inches(0.2), Inches(0.5), Inches(0.6),
                 num, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(1.4), ay + Inches(0.05), Inches(6), Inches(0.4),
                 title, size=16, color=NAVY, bold=True)
    add_rect(sl, Inches(1.4), ay + Inches(0.45), Inches(1.5), Inches(0.3), ACCENT)
    add_text_box(sl, Inches(1.45), ay + Inches(0.47), Inches(1.4), Inches(0.26),
                 when, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(1.4), ay + Inches(0.8), Inches(11), Inches(0.4),
                 desc, size=12, color=BLACK)

add_footer(sl)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"[OK] {OUT}")
print(f"     {len(prs.slides)} slides created")
